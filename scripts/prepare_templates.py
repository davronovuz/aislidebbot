"""
Slidesgo'dan yuklangan shablonlarni hybrid generator uchun tayyorlaydi:
  1. Faqat COVER (slide 1) va THANKS (oxirgi slayd) saqlanadi
  2. COVER: eng katta matn → {{TITLE}}, ikkinchisi → {{SUBTITLE}}
  3. THANKS: "Thanks!" → {{CONCLUSION_TITLE}}, kontakt placeholder'lar va
     Slidesgo watermark'i o'chiriladi
  4. Oraliq slaydlar to'liq o'chiriladi (orphan fayllar ham)
  5. manifest.json yangilanadi

Hybrid generator content slaydlarni dasturiy o'zi yasaydi.
"""

import json
from pathlib import Path

from lxml import etree
from pptx import Presentation


# ─── Konfiguratsiya ────────────────────────────────────────────────────────

TEMPLATES_DIR = Path("templates")
OUTPUT_DIR = Path("templates/prepared")

TEMPLATE_CONFIGS = {
    "Minimalist Slides.pptx": {
        "output": "minimalist.pptx",
        "name": "Minimalist",
        "description": "Toza, sodda — har mavzuga universal mos",
        "emoji": "🤍",
    },
    "Modern Education.pptx": {
        "output": "modern_edu.pptx",
        "name": "Talim",
        "description": "Talim, kurs ishi, diplom uchun",
        "emoji": "🎓",
    },
    "Multipurpose Tool.pptx": {
        "output": "multipurpose.pptx",
        "name": "Universal",
        "description": "Biznes, prezentatsiya, har mavzu",
        "emoji": "💼",
    },
}

# Watermark va placeholder matnlar — bularni topib o'chirib qo'yamiz
CLEAR_PATTERNS = [
    "youremail@",
    "yourwebsite",
    "yourname",
    "yourcompany",
    "+34 ",  # Slidesgo demo phone
    "+91 ",  # Multipurpose demo phone
    "Please keep this slide",
    "for attribution",
    "CREDITS:",
    "Slidesgo",
    "Freepik",
    "freepik.com",
    "Flaticon",
    "Storyset",
]

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ─── Yordamchilar ──────────────────────────────────────────────────────────

def should_clear(text: str) -> bool:
    t = text.strip().lower()
    for pat in CLEAR_PATTERNS:
        if pat.lower() in t:
            return True
    return False


def find_text_shapes(shape):
    """Rekursiv (group ichida ham) — barcha matnli shape'larni qaytaradi"""
    out = []
    if shape.shape_type == 6:  # GROUP
        for s in shape.shapes:
            out.extend(find_text_shapes(s))
    elif shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            try:
                fs = shape.text_frame.paragraphs[0].runs[0].font.size
                font_size = int(fs / 12700) if fs else 0
            except Exception:
                font_size = 0
            out.append({
                "shape": shape, "text": text, "fs": font_size,
                "top": getattr(shape, "top", 0) or 0,
                "left": getattr(shape, "left", 0) or 0,
            })
    return out


def set_text(text_frame, new_text: str):
    """Birinchi run formatini saqlab matnni almashtirish"""
    txBody = text_frame._txBody
    paras = txBody.findall(f"{{{NS_A}}}p")
    if not paras:
        return
    for para in paras[1:]:
        txBody.remove(para)
    first = paras[0]
    runs = first.findall(f"{{{NS_A}}}r")
    if runs:
        for r in runs[1:]:
            first.remove(r)
        t = runs[0].find(f"{{{NS_A}}}t")
        if t is not None:
            t.text = new_text
        else:
            t = etree.SubElement(runs[0], f"{{{NS_A}}}t")
            t.text = new_text
    else:
        run_elem = etree.SubElement(first, f"{{{NS_A}}}r")
        t = etree.SubElement(run_elem, f"{{{NS_A}}}t")
        t.text = new_text


def delete_slide(prs, slide_idx: int):
    """Slaydni package'dan to'liq o'chiradi (orphan fayl qolmaydi)"""
    slide = prs.slides[slide_idx]
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype.endswith("/slide") and rel._target is slide.part:
            prs.part.drop_rel(rId)
            break
    del prs.slides._sldIdLst[slide_idx]


# ─── Cover va Thanks marker assignment ─────────────────────────────────────

def mark_cover_slide(slide):
    """
    Cover'da: eng katta matn → {{TITLE}}, keyingisi → {{SUBTITLE}}.
    Watermark va kontakt'lar topilsa — o'chiriladi.
    """
    items = []
    for shape in slide.shapes:
        items.extend(find_text_shapes(shape))

    items.sort(key=lambda x: (-x["fs"], x["top"], x["left"]))

    title_set = subtitle_set = False
    for item in items:
        if should_clear(item["text"]):
            set_text(item["shape"].text_frame, "")
            continue

        if not title_set:
            set_text(item["shape"].text_frame, "{{TITLE}}")
            title_set = True
        elif not subtitle_set:
            set_text(item["shape"].text_frame, "{{SUBTITLE}}")
            subtitle_set = True
        # qolgan matnlarni o'chirib qo'yamiz (faqat title + subtitle qoladi)
        # AGAR mantiqan ko'rinish elementi bo'lsa (decorative number kabi),
        # uni saqlasak ham bo'ladi — hozir o'chiramiz toza chiqishi uchun
        else:
            # Ehtiyot: faqat aniq watermark/placeholder bo'lsa o'chir
            txt = item["text"]
            if any(c in txt for c in ["@", ".com", "+", "www.", "http"]):
                set_text(item["shape"].text_frame, "")

    return {"TITLE_SET": title_set, "SUBTITLE_SET": subtitle_set}


def mark_thanks_slide(slide):
    """
    Thanks'da: eng katta "Thanks!" → {{CONCLUSION_TITLE}}.
    "Do you have any questions?" → {{CONCLUSION_SUBTITLE}}.
    Kontakt placeholder, watermark, "Please keep..." — o'chiriladi.
    """
    items = []
    for shape in slide.shapes:
        items.extend(find_text_shapes(shape))

    items.sort(key=lambda x: (-x["fs"], x["top"], x["left"]))

    title_set = subtitle_set = False
    for item in items:
        if should_clear(item["text"]):
            set_text(item["shape"].text_frame, "")
            continue

        # @, .com, telefon raqami kabilar — o'chir
        txt = item["text"]
        if any(c in txt for c in ["@", "+34", "+91", "+1 ", "+998", "www.", "http", ".com"]):
            set_text(item["shape"].text_frame, "")
            continue

        if not title_set:
            set_text(item["shape"].text_frame, "{{CONCLUSION_TITLE}}")
            title_set = True
        elif not subtitle_set:
            set_text(item["shape"].text_frame, "{{CONCLUSION_SUBTITLE}}")
            subtitle_set = True
        else:
            # Qolgan barcha matnlarni o'chir (faqat title + subtitle qoladi)
            set_text(item["shape"].text_frame, "")

    return {"TITLE_SET": title_set, "SUBTITLE_SET": subtitle_set}


# ─── Asosiy ishlash ────────────────────────────────────────────────────────

def clear_watermarks_in_layouts(prs):
    """
    Slide layout va master'lardagi quyidagilarni tozalash:
      1. Slidesgo/Freepik watermark
      2. Default placeholder text ("Click to edit the title", "Подзаголовок"
         kabi PowerPoint defolt matnlar) — slaydda placeholder bo'sh qolganda
         ko'rinadi
    """
    DEFAULT_PLACEHOLDER_PATTERNS = [
        "click to edit",      # ENG ko'p tarqalgan default
        "click to add",
        "tap to add",
        "outline level",
        "подзагол",            # Ruscha "Подзаголовок"
        "заголовок",           # Ruscha "Заголовок"
        "title text format",
        "subtitle text",
        "edit master",
    ]

    cleaned = 0
    targets = []
    for layout in prs.slide_layouts:
        targets.append(("layout", layout))
    for master in prs.slide_masters:
        targets.append(("master", master))

    for kind, container in targets:
        for shape in container.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text
            text_lower = full_text.lower()
            # Watermark
            if any(p.lower() in text_lower for p in CLEAR_PATTERNS):
                set_text(shape.text_frame, "")
                cleaned += 1
                continue
            # Default placeholder
            if any(p in text_lower for p in DEFAULT_PLACEHOLDER_PATTERNS):
                set_text(shape.text_frame, "")
                cleaned += 1
    return cleaned


def process_template(source_name: str, config: dict) -> dict:
    src = TEMPLATES_DIR / source_name
    out = OUTPUT_DIR / config["output"]

    print(f"\n{'='*55}")
    print(f"Processing: {source_name}")

    prs = Presentation(str(src))
    n_total = len(prs.slides)
    print(f"  Slaydlar: {n_total}")

    # 1. Oraliq slaydlarni o'chirish (1, 2, ..., n-2) — teskari tartibda
    for i in range(n_total - 2, 0, -1):
        delete_slide(prs, i)
    print(f"  Oraliq o'chirildi → {len(prs.slides)} slayd qoldi (cover + thanks)")

    # 2. Layout/master'lardagi watermark'larni tozalash (CRITICAL!)
    cleaned = clear_watermarks_in_layouts(prs)
    print(f"  Layout/master watermark tozalandi: {cleaned} shape")

    # 3. Cover marker'larini qo'yish
    cover_info = mark_cover_slide(prs.slides[0])
    print(f"  Cover: TITLE={cover_info['TITLE_SET']}, SUBTITLE={cover_info['SUBTITLE_SET']}")

    # 4. Thanks marker'larini qo'yish
    thanks_info = mark_thanks_slide(prs.slides[1])
    print(f"  Thanks: TITLE={thanks_info['TITLE_SET']}, SUBTITLE={thanks_info['SUBTITLE_SET']}")

    # 4. Saqlash
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    size_kb = out.stat().st_size // 1024
    print(f"  Saqlandi: {out} ({size_kb} KB)")

    return {
        "file": config["output"],
        "name": config["name"],
        "description": config["description"],
        "emoji": config["emoji"],
        "slide_count": 2,
        "slides": [
            {"index": 0, "type": "title", "markers": ["{{TITLE}}", "{{SUBTITLE}}"]},
            {"index": 1, "type": "conclusion", "markers": ["{{CONCLUSION_TITLE}}", "{{CONCLUSION_SUBTITLE}}"]},
        ],
    }


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    manifest = {"templates": []}

    for source, config in TEMPLATE_CONFIGS.items():
        try:
            info = process_template(source, config)
            manifest["templates"].append(info)
        except Exception as e:
            print(f"XATO ({source}): {e}")
            import traceback
            traceback.print_exc()

    manifest_path = OUTPUT_DIR / "manifest.json"
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print(f"\n📋 Manifest saqlandi: {manifest_path}")
    print(f"✅ {len(manifest['templates'])} ta shablon tayyor!")


if __name__ == "__main__":
    main()
