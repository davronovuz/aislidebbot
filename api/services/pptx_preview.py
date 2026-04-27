"""Generate slide preview images and extract text from a PPTX file.

Pipeline:
  1. soffice converts .pptx → .pdf (1 PDF, all slides)
  2. pdftoppm converts each PDF page → .png
  3. python-pptx parses .pptx for text content per slide
"""
from __future__ import annotations
import json
import logging
import shutil
import subprocess
from pathlib import Path
from typing import List, Dict

logger = logging.getLogger(__name__)

DATA_ROOT = Path("/app/data/templates")


def template_dir(template_id: int) -> Path:
    d = DATA_ROOT / str(template_id)
    d.mkdir(parents=True, exist_ok=True)
    return d


def save_pptx(template_id: int, file_bytes: bytes) -> Path:
    d = template_dir(template_id)
    pptx_path = d / "template.pptx"
    pptx_path.write_bytes(file_bytes)
    return pptx_path


def _run(cmd: list[str], cwd: Path | None = None, timeout: int = 120) -> tuple[int, str, str]:
    proc = subprocess.run(
        cmd, cwd=str(cwd) if cwd else None,
        capture_output=True, text=True, timeout=timeout,
    )
    return proc.returncode, proc.stdout, proc.stderr


def generate_previews(template_id: int) -> List[str]:
    """Convert PPTX → PDF → per-slide PNG. Returns list of relative png paths."""
    d = template_dir(template_id)
    pptx_path = d / "template.pptx"
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    # Clean old previews
    for old in d.glob("slide_*.png"):
        old.unlink()
    for old in d.glob("template.pdf"):
        old.unlink()

    # Step 1: PPTX → PDF
    rc, out, err = _run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(d), str(pptx_path)],
        timeout=180,
    )
    if rc != 0:
        logger.error(f"soffice failed (rc={rc}): {err}")
        raise RuntimeError(f"PDF conversion failed: {err[:200]}")

    pdf_path = d / "template.pdf"
    if not pdf_path.exists():
        raise RuntimeError("PDF not produced")

    # Step 2: PDF → PNG per page (slide_1.png, slide_2.png, ...)
    rc, out, err = _run(
        ["pdftoppm", "-png", "-r", "120", str(pdf_path), str(d / "slide")],
        timeout=120,
    )
    if rc != 0:
        logger.error(f"pdftoppm failed: {err}")
        raise RuntimeError(f"PNG conversion failed: {err[:200]}")

    # pdftoppm names files like slide-1.png, slide-2.png — rename to slide_1.png
    pngs = sorted(d.glob("slide-*.png"))
    out_paths = []
    for i, p in enumerate(pngs, 1):
        target = d / f"slide_{i}.png"
        p.rename(target)
        out_paths.append(target.name)

    # Cleanup PDF
    pdf_path.unlink(missing_ok=True)

    return out_paths


def extract_slides_text(template_id: int) -> List[Dict]:
    """Extract title + body text from each slide using python-pptx."""
    from pptx import Presentation

    d = template_dir(template_id)
    pptx_path = d / "template.pptx"
    if not pptx_path.exists():
        return []

    pres = Presentation(str(pptx_path))
    slides = []
    for idx, slide in enumerate(pres.slides, 1):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # First non-empty text on a slide is treated as title heuristically
            if not title and shape.placeholder_format and shape.placeholder_format.idx == 0:
                title = text
            else:
                body_parts.append(text)
        if not title and body_parts:
            title = body_parts[0]
            body_parts = body_parts[1:]
        slides.append({
            "n": idx,
            "title": title,
            "body": "\n".join(body_parts),
        })

    # Save to JSON
    (d / "slides.json").write_text(json.dumps(slides, ensure_ascii=False), encoding="utf-8")
    return slides


def get_slides_data(template_id: int) -> List[Dict]:
    """Read cached slides.json (extracted text)."""
    d = template_dir(template_id)
    p = d / "slides.json"
    if not p.exists():
        return []
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return []


def list_preview_files(template_id: int) -> List[str]:
    d = template_dir(template_id)
    return sorted([p.name for p in d.glob("slide_*.png")],
                  key=lambda n: int(n.replace("slide_", "").replace(".png", "")))


def get_preview_path(template_id: int, slide_num: int) -> Path | None:
    d = template_dir(template_id)
    p = d / f"slide_{slide_num}.png"
    return p if p.exists() else None


def get_pptx_path(template_id: int) -> Path | None:
    d = template_dir(template_id)
    p = d / "template.pptx"
    return p if p.exists() else None


def delete_template_files(template_id: int) -> None:
    d = template_dir(template_id)
    if d.exists():
        shutil.rmtree(d)
