"""
HybridPPTXGenerator — gibrid yondashuv:

1. Shabloning COVER (slide 1) va THANK YOU (oxirgi slayd) saqlanadi
   — bu user tanlagan dizaynning brendigini ushlab turadi.

2. CONTENT slaydlar dasturiy yasaladi — shabloning rang palitrasidan foydalanib,
   sodda lekin chiroyli dizayn: title, bullet'lar, mavzuga mos rasm (Wikimedia).
   Bu universal har qanday mavzu uchun mos keladi.

Natija: Cover (shablon) → N ta toza content slayd → Thank You (shablon)
"""

from __future__ import annotations

import asyncio
import copy
import io
import logging
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from utils.image_provider import ImageProvider

logger = logging.getLogger(__name__)


# ─── Shablon palitralari ───────────────────────────────────────────────────
# Har shabloning cover dizaynidan ekstrapolyatsiya qilingan asosiy ranglar.
# Content slaydlar shu palitra ichida quriladi.

PALETTES: dict[str, dict[str, tuple]] = {
    "minimalist.pptx": {
        # Oq fon, sferalar ko'k-binafsha; matn qora; aksent toza ko'k
        "primary":   (15, 23, 42),       # qora matn
        "accent":    (79, 70, 229),      # indigo
        "bg":        (255, 255, 255),
        "card":      (243, 244, 246),
        "text":      (15, 23, 42),
        "muted":     (107, 114, 128),
    },
    "modern_edu.pptx": {
        # Cover'da navy + ko'k circle, kulrang ohanglar
        "primary":   (15, 35, 78),       # deep navy
        "accent":    (59, 130, 246),     # blue-500
        "bg":        (255, 255, 255),
        "card":      (240, 244, 251),
        "text":      (15, 23, 42),
        "muted":     (100, 116, 139),
    },
    "multipurpose.pptx": {
        # Cover'da to'liq ko'k fon, oq matn — content slaydlar oq + ko'k
        "primary":   (54, 62, 165),      # blue-700 (cover bilan mos)
        "accent":    (251, 146, 60),     # orange (kontrast accent)
        "bg":        (255, 255, 255),
        "card":      (239, 246, 255),
        "text":      (15, 23, 42),
        "muted":     (100, 116, 139),
    },
}

DEFAULT_PALETTE = PALETTES["minimalist.pptx"]

TEMPLATES_DIR = Path(__file__).parent.parent / "templates" / "prepared"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ─── Slayd manipulyatsiya yordamchilari (template_injector dan) ────────────


def _delete_slide(prs, slide_idx: int):
    slide = prs.slides[slide_idx]
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype.endswith("/slide") and rel._target is slide.part:
            prs.part.drop_rel(rId)
            break
    del prs.slides._sldIdLst[slide_idx]


def _reorder_slides(prs, new_order: list[int]):
    """new_order — joriy indekslar ro'yxati, target tartibi bo'yicha"""
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    for elem in list(sldIdLst):
        sldIdLst.remove(elem)
    for idx in new_order:
        sldIdLst.append(elems[idx])


# ─── Asosiy generator ──────────────────────────────────────────────────────


class HybridPPTXGenerator:
    """
    Shabloning cover + thank-you ni saqlab, oraliqda toza content slaydlar yaratuvchi
    PPTX generator. Mavzuga mos rasmlar Wikimedia/DDG dan olinadi.
    """

    def __init__(self, template_file: str):
        self.template_file = template_file
        self.palette = PALETTES.get(template_file, DEFAULT_PALETTE)

    async def generate(self, content: dict, output_path: Optional[str] = None) -> bytes:
        """
        Args:
            content: {title, subtitle, slides: [{title, content, bullet_points, image_keywords}]}
            output_path: agar berilsa, faylga ham saqlaydi

        Returns:
            PPTX bytes
        """
        template_path = TEMPLATES_DIR / self.template_file
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        prs = Presentation(str(template_path))
        slides_data = content.get("slides", [])

        # ── 1. Cover va thank-you saqlash, qolgan oraliq slaydlarni o'chirish
        n_total = len(prs.slides)
        if n_total < 2:
            raise ValueError("Shablon kamida 2 slayddan iborat bo'lishi kerak (cover + thanks)")

        # Oraliqdagi slaydlarni o'chiramiz (1, 2, 3, 4) — teskari tartibda
        for i in range(n_total - 2, 0, -1):
            _delete_slide(prs, i)
        # Endi: [cover, thank_you]

        # python-pptx ning ichki cache'sini tozalash uchun save+reload —
        # aks holda yangi slaydlar qo'shganda slide_part nomlari to'qnashadi
        _buf = io.BytesIO()
        prs.save(_buf)
        _buf.seek(0)
        prs = Presentation(_buf)

        # ── 2. Cover slaydining title va subtitle markerlarini to'ldirish
        self._fill_marker_slide(prs.slides[0], {
            "TITLE": content.get("title", "Prezentatsiya"),
            "SUBTITLE": content.get("subtitle", ""),
        })

        # ── 3. Thank you slaydining markerlarini to'ldirish
        self._fill_marker_slide(prs.slides[-1], {
            "CONCLUSION_TITLE": "Rahmat!",
            "CONCLUSION_SUBTITLE": content.get("title", ""),
        })

        # ── 4. Mavzuga mos rasmlarni parallel yuklash
        images: dict[int, str] = {}
        if slides_data:
            try:
                images = await self._fetch_images(slides_data, content.get("title", ""))
            except Exception as e:
                logger.warning(f"Rasm yuklash xato: {e}")

        # ── 5. Har content slayd uchun yangi slayd yaratish (oxiriga qo'shiladi)
        cover_layout = prs.slide_layouts[0]
        for i, slide_data in enumerate(slides_data):
            new_slide = prs.slides.add_slide(cover_layout)
            self._build_content_slide(new_slide, slide_data, images.get(i), prs.slide_width, prs.slide_height)

        # ── 6. Slaydlarni qayta tartiblash: [cover, content_1..N, thank_you]
        # Hozir: [cover(0), thanks(1), content_1(2), ..., content_N(N+1)]
        n_content = len(slides_data)
        new_order = [0] + list(range(2, 2 + n_content)) + [1]
        _reorder_slides(prs, new_order)

        # ── 7. Saqlash
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        data = buf.getvalue()

        if output_path:
            Path(output_path).write_bytes(data)

        return data

    # ── Cover/Thanks marker'larini to'ldirish ─────────────────────────────

    def _fill_marker_slide(self, slide, fill_data: dict):
        for shape in slide.shapes:
            self._fill_marker_shape(shape, fill_data)

    def _fill_marker_shape(self, shape, fill_data: dict):
        if shape.shape_type == 6:  # GROUP
            for s in shape.shapes:
                self._fill_marker_shape(s, fill_data)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith("{{") and text.endswith("}}"):
                key = text[2:-2]
                value = fill_data.get(key, "")
                self._set_text(shape.text_frame, str(value) if value else "")

    @staticmethod
    def _set_text(text_frame, new_text: str):
        txBody = text_frame._txBody
        paras = txBody.findall(f"{{{NS_A}}}p")
        if not paras:
            return
        for para in paras[1:]:
            txBody.remove(para)
        first_para = paras[0]
        runs = first_para.findall(f"{{{NS_A}}}r")
        if runs:
            for run in runs[1:]:
                first_para.remove(run)
            t_elem = runs[0].find(f"{{{NS_A}}}t")
            if t_elem is not None:
                t_elem.text = new_text
            else:
                t_elem = etree.SubElement(runs[0], f"{{{NS_A}}}t")
                t_elem.text = new_text
        else:
            run_elem = etree.SubElement(first_para, f"{{{NS_A}}}r")
            t_elem = etree.SubElement(run_elem, f"{{{NS_A}}}t")
            t_elem.text = new_text

    # ── Content slayd yasash ──────────────────────────────────────────────

    def _build_content_slide(self, slide, data: dict, image_path: Optional[str],
                              slide_w: int, slide_h: int):
        """
        Toza content slayd — slaydning haqiqiy o'lchamiga proporsional joylashadi.
        13.3x7.5 va 10x5.6 — ikkala formatda ham to'g'ri ko'rinadi.
        """
        # 1. Layoutdan kelgan barcha shape'larni o'chirish
        sp_tree = slide.shapes._spTree
        for sp_elem in list(sp_tree):
            tag = etree.QName(sp_elem).localname
            if tag in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
                sp_tree.remove(sp_elem)

        # Slaydning haqiqiy o'lchamlari (EMU). Foiz hisobida joylashtiramiz.
        W, H = slide_w, slide_h

        def pct_w(p): return int(W * p)
        def pct_h(p): return int(H * p)

        margin_x = pct_w(0.05)        # chap-o'ng 5%
        banner_h = pct_h(0.18)         # tepa banner — slayd 18%
        stripe_h = pct_h(0.012)         # accent chiziq

        # 2. Oq fon (to'liq slayd)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*self.palette["bg"])
        bg.line.fill.background()
        bg.shadow.inherit = False

        # 3. Tepa banner (primary rang)
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, banner_h)
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(*self.palette["primary"])
        banner.line.fill.background()
        banner.shadow.inherit = False

        # 4. Accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, banner_h, W, stripe_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(*self.palette["accent"])
        stripe.line.fill.background()
        stripe.shadow.inherit = False

        # 5. Title — banner ichida, oq matn
        title_top = pct_h(0.025)
        title_h = banner_h - pct_h(0.04)
        title_box = slide.shapes.add_textbox(margin_x, title_top, W - 2*margin_x, title_h)
        title_tf = title_box.text_frame
        for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(title_tf, attr, Pt(0))
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = (data.get("title") or "").strip()
        # Slayd o'lchamiga qarab title font o'lchami (proporsional)
        title_pt = 28 if W < Emu(11 * 914400) else 32
        run.font.size = Pt(title_pt)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.name = "Calibri"

        # 6. Bullet'lar va rasm bo'limi — slayd qolgan qismida
        body_top = banner_h + stripe_h + pct_h(0.06)
        body_h = H - body_top - pct_h(0.05)
        has_image = bool(image_path)

        # Bullet'lar (chap qism)
        if has_image:
            bullets_w = pct_w(0.55)
        else:
            bullets_w = W - 2*margin_x

        bullet_box = slide.shapes.add_textbox(margin_x, body_top, bullets_w, body_h)
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(btf, attr, Pt(0))

        subtitle = (data.get("content") or "").strip()
        bullets = [b for b in (data.get("bullet_points") or []) if b]

        # Slayd o'lchamiga qarab font o'lchamlari
        small = W < Emu(11 * 914400)
        sub_pt = 13 if small else 16
        bul_pt = 15 if small else 18
        dot_pt = 17 if small else 20

        first = True
        if subtitle and len(subtitle) < 200:
            p = btf.paragraphs[0]
            r = p.add_run()
            r.text = subtitle
            r.font.size = Pt(sub_pt)
            r.font.italic = True
            r.font.color.rgb = RGBColor(*self.palette["muted"])
            r.font.name = "Calibri"
            p.space_after = Pt(12)
            first = False

        for b in bullets:
            if first:
                p = btf.paragraphs[0]
                first = False
            else:
                p = btf.add_paragraph()
            r1 = p.add_run()
            r1.text = "● "
            r1.font.size = Pt(dot_pt)
            r1.font.bold = True
            r1.font.color.rgb = RGBColor(*self.palette["accent"])
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = b.strip()
            r2.font.size = Pt(bul_pt)
            r2.font.color.rgb = RGBColor(*self.palette["text"])
            r2.font.name = "Calibri"
            p.space_after = Pt(8)

        # 7. Rasm (o'ng qism) — proporsional o'lcham
        if has_image:
            try:
                img_w = pct_w(0.32)
                img_h = pct_h(0.62)
                img_left = W - margin_x - img_w
                img_top = body_top

                # Accent ramka
                frame_pad = pct_w(0.008)
                frame = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    img_left - frame_pad, img_top - frame_pad,
                    img_w + 2*frame_pad, img_h + 2*frame_pad
                )
                frame.fill.solid()
                frame.fill.fore_color.rgb = RGBColor(*self.palette["accent"])
                frame.line.fill.background()
                frame.shadow.inherit = False

                slide.shapes.add_picture(
                    image_path, img_left, img_top, width=img_w, height=img_h
                )
            except Exception as e:
                logger.warning(f"Rasm qo'yishda xato: {e}")

    # ── Rasm yuklash ──────────────────────────────────────────────────────

    async def _fetch_images(self, slides_data: list, topic: str) -> dict:
        results: dict[int, str] = {}
        async with ImageProvider() as provider:
            tasks = []
            for i, sd in enumerate(slides_data):
                title = sd.get("title", "")
                bullets = sd.get("bullet_points", []) or []
                keywords = sd.get("image_keywords", {}) or {}
                tasks.append(self._fetch_one(provider, i, title, bullets, keywords, topic))

            done = await asyncio.gather(*tasks, return_exceptions=True)
            for r in done:
                if isinstance(r, tuple):
                    idx, path = r
                    if path:
                        results[idx] = path
        return results

    @staticmethod
    async def _fetch_one(provider, idx, title, bullets, keywords, topic):
        # Avval AI bergan keyword'larni urinib ko'ramiz (eng aniq)
        for k in ("primary", "secondary"):
            kw = keywords.get(k, "")
            if kw:
                p = await provider.fetch(kw)
                if p:
                    return (idx, p)
        # Slide title + topic
        p = await provider.fetch_for_slide(title, bullets, topic)
        if p:
            return (idx, p)
        # Fallback keyword
        fb = keywords.get("fallback", "")
        if fb:
            p = await provider.fetch(fb)
            if p:
                return (idx, p)
        return (idx, None)
