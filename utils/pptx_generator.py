"""
Professional PPTX Generator — Gamma-darajaga yaqin sifat
python-pptx kutubxonasi asosida

Xususiyatlar:
- 11 ta professional rang temalari
- 5 xil slayd layouti (title, standard, card, accent-bar, image+text, xulosa)
- Gradient backgroundlar
- Shadow effektlar
- Smart text autofit
- Pixabay rasm integratsiyasi
- 16:9 format
"""

import logging
import os
import asyncio
import aiohttp
import tempfile
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)

# =====================================================================
#  RANGLAR TEMALARI — har biri to'liq rang sxemasi
# =====================================================================

THEMES = {
    "modern_blue": {
        "name": "Modern Blue",
        "title_bg": ((6, 44, 110), (14, 77, 164)),
        "slide_bg": (248, 250, 252),
        "accent": (14, 107, 247),
        "accent2": (56, 189, 248),
        "title_text": (255, 255, 255),
        "title_on_light": (10, 31, 68),
        "body_text": (51, 64, 85),
        "bullet_accent": (14, 107, 247),
        "subtitle_text": (180, 200, 230),
        "card_bg": (255, 255, 255),
        "card_border": (226, 232, 240),
    },
    "dark_elegant": {
        "name": "Dark Elegant",
        "title_bg": ((18, 18, 25), (40, 40, 60)),
        "slide_bg": (24, 24, 32),
        "accent": (139, 92, 246),
        "accent2": (167, 139, 250),
        "title_text": (255, 255, 255),
        "title_on_light": (240, 240, 255),
        "body_text": (200, 200, 215),
        "bullet_accent": (139, 92, 246),
        "subtitle_text": (150, 150, 180),
        "card_bg": (36, 36, 52),
        "card_border": (55, 55, 75),
    },
    "minimalist": {
        "name": "Minimalist",
        "title_bg": ((245, 245, 245), (235, 235, 235)),
        "slide_bg": (255, 255, 255),
        "accent": (40, 40, 40),
        "accent2": (120, 120, 120),
        "title_text": (25, 25, 25),
        "title_on_light": (25, 25, 25),
        "body_text": (60, 60, 60),
        "bullet_accent": (40, 40, 40),
        "subtitle_text": (100, 100, 100),
        "card_bg": (250, 250, 250),
        "card_border": (220, 220, 220),
    },
    "ocean_fresh": {
        "name": "Ocean Fresh",
        "title_bg": ((0, 80, 100), (0, 128, 128)),
        "slide_bg": (240, 253, 250),
        "accent": (0, 150, 136),
        "accent2": (38, 198, 218),
        "title_text": (255, 255, 255),
        "title_on_light": (0, 77, 64),
        "body_text": (38, 70, 83),
        "bullet_accent": (0, 150, 136),
        "subtitle_text": (178, 223, 219),
        "card_bg": (255, 255, 255),
        "card_border": (204, 239, 233),
    },
    "purple_premium": {
        "name": "Purple Premium",
        "title_bg": ((49, 10, 101), (88, 28, 135)),
        "slide_bg": (250, 245, 255),
        "accent": (147, 51, 234),
        "accent2": (192, 132, 252),
        "title_text": (255, 255, 255),
        "title_on_light": (49, 10, 101),
        "body_text": (59, 47, 85),
        "bullet_accent": (147, 51, 234),
        "subtitle_text": (196, 181, 253),
        "card_bg": (255, 255, 255),
        "card_border": (233, 213, 255),
    },
    "coral_warm": {
        "name": "Coral Warm",
        "title_bg": ((180, 55, 50), (220, 90, 70)),
        "slide_bg": (255, 247, 245),
        "accent": (239, 68, 68),
        "accent2": (251, 146, 60),
        "title_text": (255, 255, 255),
        "title_on_light": (127, 29, 29),
        "body_text": (68, 51, 51),
        "bullet_accent": (239, 68, 68),
        "subtitle_text": (254, 202, 202),
        "card_bg": (255, 255, 255),
        "card_border": (254, 226, 226),
    },
    "rose_creative": {
        "name": "Rose Creative",
        "title_bg": ((157, 23, 77), (190, 24, 93)),
        "slide_bg": (253, 242, 248),
        "accent": (236, 72, 153),
        "accent2": (244, 114, 182),
        "title_text": (255, 255, 255),
        "title_on_light": (131, 24, 67),
        "body_text": (76, 47, 62),
        "bullet_accent": (236, 72, 153),
        "subtitle_text": (251, 207, 232),
        "card_bg": (255, 255, 255),
        "card_border": (252, 231, 243),
    },
    "colorful_bright": {
        "name": "Colorful Bright",
        "title_bg": ((37, 99, 235), (79, 70, 229)),
        "slide_bg": (248, 250, 252),
        "accent": (79, 70, 229),
        "accent2": (16, 185, 129),
        "title_text": (255, 255, 255),
        "title_on_light": (30, 58, 138),
        "body_text": (51, 65, 85),
        "bullet_accent": (79, 70, 229),
        "subtitle_text": (191, 219, 254),
        "card_bg": (255, 255, 255),
        "card_border": (224, 231, 255),
    },
    "warm_classic": {
        "name": "Warm Classic",
        "title_bg": ((120, 80, 40), (160, 110, 60)),
        "slide_bg": (253, 251, 247),
        "accent": (180, 130, 70),
        "accent2": (210, 170, 110),
        "title_text": (255, 255, 255),
        "title_on_light": (80, 50, 20),
        "body_text": (68, 55, 40),
        "bullet_accent": (180, 130, 70),
        "subtitle_text": (220, 200, 170),
        "card_bg": (255, 255, 255),
        "card_border": (237, 228, 210),
    },
    "cosmic_dark": {
        "name": "Cosmic Dark",
        "title_bg": ((10, 5, 30), (30, 15, 60)),
        "slide_bg": (16, 12, 38),
        "accent": (99, 102, 241),
        "accent2": (236, 72, 153),
        "title_text": (255, 255, 255),
        "title_on_light": (230, 230, 255),
        "body_text": (200, 200, 225),
        "bullet_accent": (99, 102, 241),
        "subtitle_text": (150, 150, 200),
        "card_bg": (28, 22, 56),
        "card_border": (50, 42, 80),
    },
    "green_nature": {
        "name": "Green Nature",
        "title_bg": ((6, 78, 59), (21, 128, 61)),
        "slide_bg": (240, 253, 244),
        "accent": (22, 163, 74),
        "accent2": (74, 222, 128),
        "title_text": (255, 255, 255),
        "title_on_light": (6, 78, 59),
        "body_text": (38, 70, 50),
        "bullet_accent": (22, 163, 74),
        "subtitle_text": (187, 247, 208),
        "card_bg": (255, 255, 255),
        "card_border": (209, 250, 229),
    },
}

# Mavjud bot theme ID → yangi tema mapping
THEME_ID_MAP = {
    "chisel": "minimalist",
    "coal": "dark_elegant",
    "blues": "modern_blue",
    "elysia": "rose_creative",
    "breeze": "ocean_fresh",
    "aurora": "purple_premium",
    "coral-glow": "coral_warm",
    "gamma": "colorful_bright",
    "creme": "warm_classic",
    "gamma-dark": "cosmic_dark",
}

DEFAULT_THEME = "modern_blue"

# Slayd o'lchamlari (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =====================================================================
#  GENERATOR CLASS
# =====================================================================

class ProPPTXGenerator:
    """
    Professional PPTX generator — python-pptx asosida

    Foydalanish:
        gen = ProPPTXGenerator(theme_id="blues")
        success = await gen.generate(content, "output.pptx")
    """

    def __init__(self, theme_id: str = None):
        theme_name = DEFAULT_THEME
        if theme_id:
            theme_name = THEME_ID_MAP.get(theme_id.lower(), theme_id.lower())
        if theme_name not in THEMES:
            theme_name = DEFAULT_THEME

        self.theme = THEMES[theme_name]
        self.theme_name = theme_name
        self.prs = None

    # ======================== MAIN API ========================

    async def generate(
        self,
        content: Dict,
        output_path: str,
        pixabay_api_key: str = None,
    ) -> bool:
        """
        Professional PPTX yaratish

        Args:
            content: GPT-4o dan kelgan content dict
            output_path: Chiqish fayl yo'li
            pixabay_api_key: Pixabay API kaliti (ixtiyoriy)

        Returns:
            True — muvaffaqiyatli
        """
        try:
            # 1. Rasmlarni yuklab olish (agar kalit berilgan bo'lsa)
            images = {}
            if pixabay_api_key:
                images = await self._fetch_images(content, pixabay_api_key)

            # 2. PPTX yaratish
            self._build(content, images, output_path)

            # 3. Vaqtinchalik rasmlarni tozalash
            for img_path in images.values():
                try:
                    if img_path and os.path.exists(img_path):
                        os.remove(img_path)
                except Exception:
                    pass

            file_size = os.path.getsize(output_path)
            logger.info(f"PPTX yaratildi: {output_path} ({file_size:,} bytes, theme: {self.theme_name})")
            return True

        except Exception as e:
            logger.error(f"PPTX generate xato: {e}", exc_info=True)
            return False

    # ======================== BUILD ========================

    def _build(self, content: Dict, images: Dict, output_path: str):
        """PPTX ni qurib saqlash"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

        title = content.get("title", "Prezentatsiya")
        subtitle = content.get("subtitle", "")
        slides = content.get("slides", [])

        # Title slayd
        self._create_title_slide(title, subtitle)

        # Content slaydlar
        for i, slide_data in enumerate(slides):
            is_last = (i == len(slides) - 1)

            if is_last and len(slides) > 1:
                self._create_conclusion_slide(slide_data)
            else:
                variant = i % 3  # 0=standard, 1=card, 2=accent-bar
                img_path = images.get(i)
                self._create_content_slide(slide_data, variant, img_path)

        self.prs.save(output_path)

    # ======================== TITLE SLIDE ========================

    def _create_title_slide(self, title: str, subtitle: str):
        """Chiroyli title slayd — gradient bg, accent elementlar"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        t = self.theme

        # Gradient background
        self._set_gradient_bg(slide, t["title_bg"][0], t["title_bg"][1])

        # Yuqori dekorativ chiziq
        self._add_rect(slide, Inches(1), Inches(0.5),
                        Inches(2.5), Inches(0.07), t["accent2"])

        # Kichik dekorativ kvadrat
        self._add_rect(slide, Inches(11.5), Inches(0.5),
                        Inches(0.5), Inches(0.5), t["accent"], alpha=40)

        # Title
        self._add_textbox(
            slide, title,
            x=Inches(1), y=Inches(2.0),
            w=Inches(11.333), h=Inches(2.2),
            font_size=44, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Accent chiziq title ostida
        self._add_rect(slide, Inches(1), Inches(4.3),
                        Inches(4), Inches(0.08), t["accent"])

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                x=Inches(1), y=Inches(4.65),
                w=Inches(9), h=Inches(1.3),
                font_size=20, bold=False,
                color=t["subtitle_text"],
                alignment=PP_ALIGN.LEFT,
            )

        # Pastki gradient chiziq
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.12),
                        SLIDE_W, Inches(0.12), t["accent"])

        # Pastki o'ng burchak dekorativ element
        self._add_rect(slide, SLIDE_W - Inches(3), SLIDE_H - Inches(0.12),
                        Inches(3), Inches(0.12), t["accent2"])

    # ======================== CONTENT SLIDES ========================

    def _create_content_slide(self, data: Dict, variant: int, image_path: str = None):
        """Content slayd — layout variant tanlash"""
        if image_path and os.path.exists(image_path) and variant != 2:
            self._create_image_content_slide(data, image_path)
        elif variant == 0:
            self._create_standard_slide(data)
        elif variant == 1:
            self._create_card_slide(data)
        else:
            self._create_accent_bar_slide(data)

    def _create_standard_slide(self, data: Dict):
        """Standard layout — rangli title bar yuqorida, kontent pastda"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title bar (to'liq kenglik)
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.5), t["title_bg"][0])

        # Title bar ostida accent chiziq
        self._add_rect(slide, Inches(0), Inches(1.5),
                        SLIDE_W, Inches(0.06), t["accent"])

        # Title text
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.3),
            w=Inches(11.7), h=Inches(0.9),
            font_size=30, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Content va bullets
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(2.0)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(0.8), y=y,
                w=Inches(11.7), h=Inches(1.6),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.7)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(1.0), y=y,
                w=Inches(11.3), h=SLIDE_H - y - Inches(0.5),
                font_size=16, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

        # Pastki chiziq
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.06),
                        SLIDE_W, Inches(0.06), t["accent"])

    def _create_card_slide(self, data: Dict):
        """Card layout — kontent karta ichida, soyali"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title (karta ustida)
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.4),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_on_light"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Title ostida accent chiziq
        self._add_rect(slide, Inches(0.8), Inches(1.3),
                        Inches(3), Inches(0.06), t["accent"])

        # Card — rounded rect, shadow bilan
        card = self._add_rounded_rect(
            slide,
            x=Inches(0.6), y=Inches(1.75),
            w=Inches(12.1), h=Inches(5.2),
            fill=t["card_bg"],
            border_color=t.get("card_border"),
            shadow=True,
        )

        # Kartaning chap tomonida accent chiziq
        self._add_rect(slide, Inches(0.6), Inches(1.75),
                        Inches(0.07), Inches(5.2), t["accent"])

        # Card ichida content
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(2.15)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(1.3), y=y,
                w=Inches(10.7), h=Inches(1.5),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.7)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(1.5), y=y,
                w=Inches(10.2), h=Inches(3.0),
                font_size=16, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

    def _create_accent_bar_slide(self, data: Dict):
        """Accent bar layout — chapda rangli bar, ikki ustunli bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Chap katta accent bar
        self._add_rect(slide, Inches(0), Inches(0),
                        Inches(0.45), SLIDE_H, t["title_bg"][0])

        # Bar yonida ingichka accent chiziq
        self._add_rect(slide, Inches(0.45), Inches(0),
                        Inches(0.06), SLIDE_H, t["accent"])

        # Yuqori o'ng burchakda kichik dekorativ element
        self._add_rect(slide, SLIDE_W - Inches(2), Inches(0.4),
                        Inches(1.2), Inches(0.06), t["accent2"])

        # Title
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(1.1), y=Inches(0.45),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_on_light"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Divider
        self._add_rect(slide, Inches(1.1), Inches(1.4),
                        Inches(2.5), Inches(0.05), t["accent"])

        # Content
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(1.85)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(1.1), y=y,
                w=Inches(11.7), h=Inches(1.5),
                font_size=17, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(3.5)

        # Bullets — 4+ bo'lsa ikki ustunli
        if bullets:
            if len(bullets) >= 4:
                mid = (len(bullets) + 1) // 2
                self._add_bullet_textbox(
                    slide, bullets[:mid],
                    x=Inches(1.1), y=y,
                    w=Inches(5.5), h=SLIDE_H - y - Inches(0.5),
                    font_size=15, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )
                self._add_bullet_textbox(
                    slide, bullets[mid:],
                    x=Inches(7.0), y=y,
                    w=Inches(5.5), h=SLIDE_H - y - Inches(0.5),
                    font_size=15, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )
            else:
                self._add_bullet_textbox(
                    slide, bullets,
                    x=Inches(1.1), y=y,
                    w=Inches(11.7), h=SLIDE_H - y - Inches(0.5),
                    font_size=16, color=t["body_text"],
                    bullet_color=t["bullet_accent"],
                )

    def _create_image_content_slide(self, data: Dict, image_path: str):
        """Rasm + kontent layout — chapda rasm, o'ngda matn"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        self._set_solid_bg(slide, t["slide_bg"])

        # Title bar
        self._add_rect(slide, Inches(0), Inches(0),
                        SLIDE_W, Inches(1.4), t["title_bg"][0])
        self._add_rect(slide, Inches(0), Inches(1.4),
                        SLIDE_W, Inches(0.05), t["accent"])

        # Title
        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            x=Inches(0.8), y=Inches(0.28),
            w=Inches(11.7), h=Inches(0.85),
            font_size=28, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.LEFT,
            font_name="Calibri Light",
        )

        # Rasm (chap tomonda, soyali frame ichida)
        img_frame = self._add_rounded_rect(
            slide,
            x=Inches(0.5), y=Inches(1.8),
            w=Inches(5.8), h=Inches(5.1),
            fill=t["card_bg"],
            border_color=t.get("card_border"),
            shadow=True,
        )

        try:
            slide.shapes.add_picture(
                image_path,
                left=Inches(0.65), top=Inches(1.95),
                width=Inches(5.5), height=Inches(4.8),
            )
        except Exception as e:
            logger.warning(f"Rasm qo'shishda xato: {e}")

        # O'ng tomonda kontent
        content = data.get("content", "")
        bullets = data.get("bullet_points", [])
        y = Inches(1.85)

        if content:
            self._add_textbox(
                slide, content,
                x=Inches(6.7), y=y,
                w=Inches(6.1), h=Inches(2.0),
                font_size=16, bold=False,
                color=t["body_text"],
                line_spacing=1.4,
            )
            y = Inches(4.0)

        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(6.9), y=y,
                w=Inches(5.9), h=SLIDE_H - y - Inches(0.5),
                font_size=15, color=t["body_text"],
                bullet_color=t["bullet_accent"],
            )

    # ======================== CONCLUSION SLIDE ========================

    def _create_conclusion_slide(self, data: Dict):
        """Xulosa slayd — gradient bg, markazlashtirilgan"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.theme

        # Gradient bg (boshqa yo'nalishda)
        self._set_gradient_bg(slide, t["title_bg"][1], t["title_bg"][0], angle=2700000)

        # Yuqori dekorativ chiziq (markazda)
        bar_w = Inches(4)
        bar_x = (SLIDE_W - bar_w) // 2
        self._add_rect(slide, bar_x, Inches(1.8),
                        bar_w, Inches(0.06), t["accent2"])

        # Title
        title = data.get("title", "Xulosa")
        self._add_textbox(
            slide, title,
            x=Inches(1), y=Inches(2.2),
            w=Inches(11.333), h=Inches(1.3),
            font_size=40, bold=True,
            color=t["title_text"],
            alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light",
        )

        # Content
        content = data.get("content", "")
        if content:
            self._add_textbox(
                slide, content,
                x=Inches(2), y=Inches(3.8),
                w=Inches(9.333), h=Inches(1.8),
                font_size=19, bold=False,
                color=t["subtitle_text"],
                alignment=PP_ALIGN.CENTER,
                line_spacing=1.5,
            )

        # Bullets (markazda)
        bullets = data.get("bullet_points", [])
        if bullets:
            self._add_bullet_textbox(
                slide, bullets,
                x=Inches(3), y=Inches(4.8),
                w=Inches(7.333), h=Inches(2.0),
                font_size=16, color=t["subtitle_text"],
                bullet_color=t["accent2"],
                alignment=PP_ALIGN.CENTER,
            )

        # Pastki accent chiziqlar
        self._add_rect(slide, Inches(0), SLIDE_H - Inches(0.12),
                        SLIDE_W, Inches(0.12), t["accent"])
        self._add_rect(slide, SLIDE_W - Inches(3.5), SLIDE_H - Inches(0.12),
                        Inches(3.5), Inches(0.12), t["accent2"])

    # ======================== TEXT HELPERS ========================

    def _add_textbox(self, slide, text: str, x, y, w, h,
                     font_size=14, bold=False, color=(0, 0, 0),
                     alignment=PP_ALIGN.LEFT, line_spacing=1.2,
                     font_name="Calibri"):
        """Text box qo'shish — multiline, auto-fit"""
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)

        lines = text.split('\n')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = font_name
            p.font.color.rgb = RGBColor(*color)
            p.alignment = alignment
            p.space_after = Pt(3)
            p.space_before = Pt(1)

            self._set_line_spacing(p, line_spacing)

        # Shrink text on overflow
        self._set_text_autofit(txBox)
        return txBox

    def _add_bullet_textbox(self, slide, bullets: List[str], x, y, w, h,
                            font_size=16, color=(0, 0, 0),
                            bullet_color=(0, 0, 0),
                            alignment=PP_ALIGN.LEFT):
        """Bullet pointlar uchun maxsus text box — har bir bullet alohida paragraf"""
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        for idx, bullet in enumerate(bullets):
            bullet = bullet.strip()
            if not bullet:
                continue

            # Mavjud bullet belgilarini olib tashlash
            for prefix in ('•', '●', '○', '▪', '▸', '-', '–', '—', '*'):
                if bullet.startswith(prefix):
                    bullet = bullet[len(prefix):].strip()
                    break

            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Bullet marker (rangli) + matn
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(font_size - 2)
            run_bullet.font.color.rgb = RGBColor(*bullet_color)
            run_bullet.font.name = "Calibri"

            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = RGBColor(*color)
            run_text.font.name = "Calibri"

            p.alignment = alignment
            p.space_after = Pt(8)
            p.space_before = Pt(4)

            self._set_line_spacing(p, 1.3)

        self._set_text_autofit(txBox)
        return txBox

    # ======================== SHAPE HELPERS ========================

    def _add_rect(self, slide, x, y, w, h, fill, alpha=None):
        """To'rtburchak shape qo'shish"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.fill.background()

        if alpha is not None:
            self._set_shape_alpha(shape, alpha)

        return shape

    def _add_rounded_rect(self, slide, x, y, w, h,
                          fill=(255, 255, 255),
                          border_color=None, shadow=False):
        """Rounded rectangle — burchaklari yumaloq"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)

        if border_color:
            shape.line.color.rgb = RGBColor(*border_color)
            shape.line.width = Pt(0.75)
        else:
            shape.line.fill.background()

        if shadow:
            self._add_shadow(shape)

        return shape

    # ======================== BACKGROUND HELPERS ========================

    def _set_gradient_bg(self, slide, color1, color2, angle=5400000):
        """Slide ga gradient background (XML orqali)"""
        bg = slide.background._element

        for child in list(bg):
            bg.remove(child)

        bgPr = etree.SubElement(bg, qn('p:bgPr'))
        gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        # Color stop 1
        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '0')
        srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
        srgb1.set('val', '%02X%02X%02X' % color1)

        # Color stop 2
        gs2 = etree.SubElement(gsLst, qn('a:gs'))
        gs2.set('pos', '100000')
        srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
        srgb2.set('val', '%02X%02X%02X' % color2)

        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(angle))
        lin.set('scaled', '1')

        etree.SubElement(bgPr, qn('a:effectLst'))

    def _set_solid_bg(self, slide, color):
        """Slide ga solid background"""
        bg = slide.background._element

        for child in list(bg):
            bg.remove(child)

        bgPr = etree.SubElement(bg, qn('p:bgPr'))
        solidFill = etree.SubElement(bgPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', '%02X%02X%02X' % color)

        etree.SubElement(bgPr, qn('a:effectLst'))

    # ======================== XML EFFECTS ========================

    def _add_shadow(self, shape):
        """Shape ga drop shadow (XML orqali)"""
        try:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is None:
                return

            effectLst = spPr.find(qn('a:effectLst'))
            if effectLst is None:
                effectLst = etree.SubElement(spPr, qn('a:effectLst'))

            outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
            outerShdw.set('blurRad', '63500')
            outerShdw.set('dist', '25400')
            outerShdw.set('dir', '5400000')
            outerShdw.set('algn', 'tl')
            outerShdw.set('rotWithShape', '0')

            srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
            srgbClr.set('val', '000000')
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '22000')
        except Exception as e:
            logger.debug(f"Shadow xato: {e}")

    def _set_shape_alpha(self, shape, alpha_pct):
        """Shape ga shaffoflik (0-100, 0=to'liq shaffof)"""
        try:
            spPr = shape._element.find(qn('p:spPr'))
            if spPr is None:
                return

            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is None:
                return

            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is None:
                return

            a = etree.SubElement(srgbClr, qn('a:alpha'))
            a.set('val', str(alpha_pct * 1000))
        except Exception as e:
            logger.debug(f"Alpha xato: {e}")

    def _set_text_autofit(self, shape):
        """Matn sig'masa avtomatik kichiklaydigan qilish"""
        try:
            txBody = shape._element.find(qn('p:txBody'))
            if txBody is None:
                return

            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is None:
                return

            bodyPr.set('wrap', 'square')
            bodyPr.set('anchor', 't')

            for tag in ['a:noAutofit', 'a:normAutofit', 'a:spAutoFit']:
                existing = bodyPr.find(qn(tag))
                if existing is not None:
                    bodyPr.remove(existing)

            normAutofit = etree.SubElement(bodyPr, qn('a:normAutofit'))
            normAutofit.set('fontScale', '100000')
            normAutofit.set('lnSpcReduction', '20000')
        except Exception as e:
            logger.debug(f"Autofit xato: {e}")

    def _set_line_spacing(self, paragraph, multiplier: float):
        """Paragrafga qator oralig'i qo'yish"""
        try:
            pPr = paragraph._element.find(qn('a:pPr'))
            if pPr is None:
                pPr = etree.SubElement(paragraph._element, qn('a:pPr'))

            # Mavjud lnSpc ni olib tashlash
            existing = pPr.find(qn('a:lnSpc'))
            if existing is not None:
                pPr.remove(existing)

            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(multiplier * 100000)))
        except Exception as e:
            logger.debug(f"Line spacing xato: {e}")

    # ======================== IMAGE FETCHING ========================

    async def _fetch_images(self, content: Dict, api_key: str) -> Dict[int, str]:
        """Pixabay dan rasmlar yuklab olish — har bir slayd uchun"""
        images = {}
        slides = content.get("slides", [])

        timeout = aiohttp.ClientTimeout(total=30)
        async with aiohttp.ClientSession(timeout=timeout) as session:
            tasks = []
            for i, slide in enumerate(slides):
                keywords = slide.get("image_keywords", {})
                if keywords:
                    tasks.append(self._fetch_slide_image(session, api_key, i, keywords))

            results = await asyncio.gather(*tasks, return_exceptions=True)
            for result in results:
                if isinstance(result, tuple):
                    idx, path = result
                    if path:
                        images[idx] = path

        logger.info(f"{len(images)} ta rasm yuklab olindi")
        return images

    async def _fetch_slide_image(self, session, api_key: str,
                                  slide_idx: int, keywords: Dict) -> Tuple[int, Optional[str]]:
        """Bitta slayd uchun rasm yuklab olish — 3 bosqichli fallback"""
        for key_type in ["primary", "secondary", "fallback"]:
            keyword = keywords.get(key_type, "")
            if not keyword:
                continue

            img_path = await self._download_pixabay_image(session, api_key, keyword)
            if img_path:
                return (slide_idx, img_path)

        return (slide_idx, None)

    async def _download_pixabay_image(self, session, api_key: str,
                                       keyword: str) -> Optional[str]:
        """Pixabay API dan bitta rasm yuklab olish"""
        try:
            search_url = (
                f"https://pixabay.com/api/"
                f"?key={api_key}"
                f"&q={keyword}"
                f"&image_type=photo"
                f"&orientation=horizontal"
                f"&per_page=3"
                f"&min_width=800"
                f"&safesearch=true"
            )

            async with session.get(search_url) as resp:
                if resp.status != 200:
                    return None
                data = await resp.json()

            hits = data.get("hits", [])
            if not hits:
                return None

            img_url = hits[0].get("webformatURL", "")
            if not img_url:
                return None

            async with session.get(img_url) as resp:
                if resp.status != 200:
                    return None
                img_data = await resp.read()

            tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            tmp.write(img_data)
            tmp.close()
            return tmp.name

        except Exception as e:
            logger.debug(f"Pixabay xato ({keyword}): {e}")
            return None


# =====================================================================
#  YORDAMCHI FUNKSIYALAR
# =====================================================================

def get_available_themes() -> List[Dict]:
    """Mavjud temalar ro'yxati"""
    result = []
    for theme_id, theme_name in THEME_ID_MAP.items():
        theme = THEMES.get(theme_name, {})
        result.append({
            "id": theme_id,
            "internal_name": theme_name,
            "display_name": theme.get("name", theme_name),
        })
    return result


def resolve_theme_id(theme_id: str) -> str:
    """Bot theme ID ni generator theme nomiga aylantirish"""
    if not theme_id:
        return DEFAULT_THEME
    return THEME_ID_MAP.get(theme_id.lower(), DEFAULT_THEME)
