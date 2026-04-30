"""
Template Injector: tayyorlangan .pptx shablonlarni AI kontent bilan to'ldiradi.

Asosiy texnik nuqta — slayd KO'CHIRGANDA (clone) shablon ichidagi rasmlar
saqlanishi kerak. Buning uchun XML deepcopy bilan birga slide_part'ning
relationship'larini ham yangi slaydga rId remapping qilib ko'chirish kerak.
Aks holda `r:embed="rIdX"` referenсiyalari buziladi va dizayn yo'qoladi.
"""

import copy
import io
import json
import logging
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# rId saqlovchi XML attributlar — bularni remap qilish kerak
RID_ATTRS = {
    f"{{{NS_R}}}embed",   # rasm: <a:blip r:embed="rId3"/>
    f"{{{NS_R}}}link",    # tashqi rasm/link
    f"{{{NS_R}}}id",      # rels (header/footer/notes)
}

TEMPLATES_DIR = Path(__file__).parent.parent / "templates" / "prepared"


# ─── Slayd manipulyatsiya yordamchilari ────────────────────────────────────


def _copy_part_rels(source_part, target_part) -> dict:
    """
    Source part'ning barcha relationship'larini target part'ga ko'chiradi
    (slideLayout va notesSlide bundan mustasno — ular tegishlimas).
    Eski rId → yangi rId xaritasi qaytariladi.
    """
    rId_map = {}
    for rel in list(source_part.rels.values()):
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.reltype.endswith("/notesSlide"):
            continue
        try:
            if rel.is_external:
                new_rId = target_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                new_rId = target_part.relate_to(rel.target_part, rel.reltype)
            rId_map[rel.rId] = new_rId
        except Exception as e:
            logger.debug(f"Rel ko'chirib bo'lmadi {rel.rId}: {e}")
    return rId_map


def _remap_rids_in_xml(elem, rId_map: dict):
    """XML daraxtini aylanib, rId reference'larini yangi qiymatlarga almashtiradi"""
    for node in elem.iter():
        for attr in RID_ATTRS:
            if attr in node.attrib:
                old = node.attrib[attr]
                if old in rId_map:
                    node.attrib[attr] = rId_map[old]


def _clear_target_rels(target_part):
    """Target slide'ning slideLayout'dan boshqa barcha rels'ini o'chiradi"""
    for rId in list(target_part.rels.keys()):
        rel = target_part.rels[rId]
        if rel.reltype.endswith("/slideLayout"):
            continue
        try:
            target_part.drop_rel(rId)
        except Exception:
            pass


def _replace_slide_with(target_slide, source_slide):
    """
    Target slide'ni source slide'ning to'liq nusxasiga aylantiradi:
    - source ning rels'ini target'ga remap bilan ko'chiradi
    - source ning spTree XML'ini target'ga ko'chiradi (rId'lar yangilangan)
    """
    target_part = target_slide.part
    source_part = source_slide.part

    # 1. Target'ning eski rels'ini tozalash (slideLayout qoladi)
    _clear_target_rels(target_part)

    # 2. Source rels'ni target'ga ko'chirish, rId xaritasini olish
    rId_map = _copy_part_rels(source_part, target_part)

    # 3. Target spTree'ni butunlay almashtirish
    target_tree = target_slide.shapes._spTree
    for child in list(target_tree):
        target_tree.remove(child)

    # 4. Source spTree dan har bir shape ni copy qilib, rId'larini remap qilib qo'shish
    for child in source_slide.shapes._spTree:
        copied = copy.deepcopy(child)
        _remap_rids_in_xml(copied, rId_map)
        target_tree.append(copied)


def _clone_slide_at_end(prs, source_slide):
    """
    Yangi slayd qo'shadi va source_slide'ning to'liq nusxasiga aylantiradi.
    Dizayn, rasmlar, format - hammasi saqlanadi.
    """
    # Bo'sh slayd qo'shamiz (slideLayout source bilan bir xil)
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    _replace_slide_with(new_slide, source_slide)
    return new_slide


def _delete_slide(prs, slide_idx: int):
    """Slide ID list'dan va presentation rels'idan o'chirish"""
    slide_to_delete = prs.slides[slide_idx]
    # rels o'chirish
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype.endswith("/slide") and rel._target is slide_to_delete.part:
            prs.part.drop_rel(rId)
            break
    # Slide ID list'dan
    del prs.slides._sldIdLst[slide_idx]


# ─── Asosiy injector klassi ────────────────────────────────────────────────


class TemplateInjector:
    def __init__(self):
        manifest_path = TEMPLATES_DIR / "manifest.json"
        with open(manifest_path, encoding="utf-8") as f:
            self.manifest = json.load(f)
        self._templates = {t["file"]: t for t in self.manifest["templates"]}

    # ── Public API ─────────────────────────────────────────────────────────

    def get_templates(self) -> list[dict]:
        return [
            {"file": t["file"], "name": t["name"], "description": t["description"]}
            for t in self.manifest["templates"]
        ]

    def get_template_slots(self, template_file: str) -> dict:
        info = self._templates.get(template_file)
        if not info:
            return {}
        content_slides = [s for s in info["slides"] if s["type"] == "content"]
        slots = []
        for s in content_slides:
            n = sum(1 for m in s["markers"] if m.startswith("{{BULLET_"))
            slots.append(n)
        return {
            "content_slide_count": len(content_slides),
            "bullets_per_slide": slots,
            "max_bullets": max(slots) if slots else 4,
            "min_bullets": min(slots) if slots else 2,
        }

    def inject(self, template_file: str, content: dict) -> bytes:
        """
        Asosiy funksiya: shablonga AI kontentni joylashtiradi.
        Slaydlar to'liq nusxalanadi (rasmlar, icon'lar, dizayn — hammasi saqlanadi).
        """
        info = self._templates.get(template_file)
        if not info:
            raise ValueError(f"Shablon topilmadi: {template_file}")

        prs = Presentation(str(TEMPLATES_DIR / template_file))
        slides_data = content.get("slides", [])

        # Template'dagi slayd indekslari (turi bo'yicha)
        title_idx = next(i for i, s in enumerate(info["slides"]) if s["type"] == "title")
        conclusion_idx = next(i for i, s in enumerate(info["slides"]) if s["type"] == "conclusion")
        content_idxs = [i for i, s in enumerate(info["slides"]) if s["type"] == "content"]

        if not content_idxs:
            raise ValueError(f"Shablonda content slayd yo'q: {template_file}")

        # Maqsadli ketma-ketlik: [title, content×N, conclusion]
        target_seq = [title_idx]
        for i in range(len(slides_data)):
            target_seq.append(content_idxs[i % len(content_idxs)])
        target_seq.append(conclusion_idx)

        # ── BOSQICH 1: Manba slaydlarini o'zgarmas saqlaymiz ────────────
        # Original prs.slides'da har slaydning to'liq part'i (rels + spTree) bor.
        # Birinchi navbatda — barcha slaydlarning ZAXIRA ko'rsatkichlarini olamiz.
        original_slides = list(prs.slides)
        # Eslatma: bu list slide obyektlariga referensa, prs o'zgarsa ham ular qoladi

        # ── BOSQICH 2: Yangi slaydlar yaratamiz target_seq tartibida ──
        # Har bir target_seq[i] uchun, original_slides[target_seq[i]]'ni
        # to'liq clone qilib, prezentatsiya OXIRIGA qo'shamiz.
        new_slides = []
        for tpl_idx in target_seq:
            source = original_slides[tpl_idx]
            cloned = _clone_slide_at_end(prs, source)
            new_slides.append(cloned)

        # ── BOSQICH 3: Eski (asl) slaydlarni o'chiramiz ────────────────
        # Hozir: [orig_0..orig_5, new_0..new_n] — birinchi 6 tasini o'chirish
        n_original = len(original_slides)
        for _ in range(n_original):
            _delete_slide(prs, 0)

        # ── BOSQICH 4: Kontent joylash ─────────────────────────────────
        prs_title = content.get("title", "Prezentatsiya")
        prs_subtitle = content.get("subtitle", "")

        # Title slayd
        self._fill_slide(new_slides[0], {
            "TITLE": prs_title,
            "SUBTITLE": prs_subtitle,
        })

        # Content slaydlar
        for slide, slide_data in zip(new_slides[1:-1], slides_data):
            fill = {
                "TITLE": slide_data.get("title", ""),
                "SUBTITLE": slide_data.get("content", ""),
                "CONTENT": slide_data.get("content", ""),
            }
            for j, bp in enumerate(slide_data.get("bullet_points", [])[:5]):
                fill[f"BULLET_{j + 1}"] = bp
            self._fill_slide(slide, fill)

        # Xulosa slayd
        self._fill_slide(new_slides[-1], {
            "CONCLUSION_TITLE": "Xulosa",
            "CONCLUSION_SUBTITLE": prs_title,
        })

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    # ── Matn to'ldirish ───────────────────────────────────────────────────

    def _fill_slide(self, slide, fill_data: dict):
        for shape in slide.shapes:
            self._fill_shape(shape, fill_data)

    def _fill_shape(self, shape, fill_data: dict):
        if shape.shape_type == 6:  # GROUP
            for s in shape.shapes:
                self._fill_shape(s, fill_data)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text.startswith("{{") and text.endswith("}}"):
                key = text[2:-2]
                value = fill_data.get(key, "")
                self._set_text(shape.text_frame, str(value) if value else "")

    def _set_text(self, text_frame, new_text: str):
        """Matnni almashtirish — birinchi run formatini saqlash"""
        txBody = text_frame._txBody
        paras = txBody.findall(f"{{{NS_A}}}p")
        if not paras:
            return

        for para in paras[1:]:
            txBody.remove(para)

        first_para = paras[0]
        runs = first_para.findall(f"{{{NS_A}}}r")

        if runs:
            for run in runs[1:]:
                first_para.remove(run)
            t_elem = runs[0].find(f"{{{NS_A}}}t")
            if t_elem is not None:
                t_elem.text = new_text
            else:
                t_elem = etree.SubElement(runs[0], f"{{{NS_A}}}t")
                t_elem.text = new_text
        else:
            run_elem = etree.SubElement(first_para, f"{{{NS_A}}}r")
            t_elem = etree.SubElement(run_elem, f"{{{NS_A}}}t")
            t_elem.text = new_text


def get_injector() -> TemplateInjector:
    if not hasattr(get_injector, "_instance"):
        get_injector._instance = TemplateInjector()
    return get_injector._instance
