from pptx import Presentation
from pptx.util import Inches, Pt
import os

path = "/tmp/test_full.pptx"
prs = Presentation(path)
print("Fayl:", os.path.getsize(path), "bytes")
print("Slide W:", prs.slide_width, "H:", prs.slide_height)
print("Slaydlar:", len(prs.slides))
print()

for i, slide in enumerate(prs.slides):
    print("=== SLAYD", i+1, "===")
    print("  Shapes:", len(slide.shapes))

    for j, shape in enumerate(slide.shapes):
        stype = str(shape.shape_type)
        print("  Shape", j, ":", shape.name, "|", stype)
        print("    left=%d top=%d w=%d h=%d" % (shape.left, shape.top, shape.width, shape.height))
        if shape.has_text_frame:
            text = shape.text_frame.text[:100]
            print("    TEXT:", repr(text))
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    fs = str(r.font.size)
                    try:
                        fc = str(r.font.color.rgb)
                    except:
                        fc = "None"
                    print("      run: size=%s color=%s [%s]" % (fs, fc, r.text[:50]))
    print()
